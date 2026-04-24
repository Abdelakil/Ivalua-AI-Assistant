"""Incremental documentation importer for ConPort.

Scans Documentation/ for text-extractable files, compares against a manifest
of previously-imported files, and only re-processes new or changed files.
Removed files are purged from the DB. Module indexes are rebuilt on every run.

Usage:
    python context_portal/scripts/import_documentation.py

Dependencies (from pyproject.toml):
    PyPDF2, python-docx, openpyxl, python-pptx
"""
from __future__ import annotations

import json
import re
import sqlite3
import sys
from datetime import datetime, timezone
from pathlib import Path

WORKSPACE = Path(__file__).resolve().parents[2]
DB = WORKSPACE / "context_portal" / "context.db"
DOCS_ROOT = WORKSPACE / "Documentation"
MANIFEST_PATH = WORKSPACE / "context_portal" / "import_data" / "doc_manifest.json"

CHUNK_SIZE = 3000
CHUNK_OVERLAP = 500
MAX_CHUNK_BYTES = 50_000

CATEGORY = "IVALUA_Documentation"

SKIP_EXT = {".mp4", ".jpg", ".jpeg", ".png", ".mp3", ".tscproj", ".csv"}
TEXT_EXT = {".pdf", ".txt", ".docx", ".xlsx", ".pptx"}

COPYRIGHT_RE = re.compile(r"©\s*\d{4}.*Ivalua.*All rights reserved", re.IGNORECASE)
PAGE_NUM_RE = re.compile(r"^\s*(Page\s+\d+\s+of\s+\d+|\d+\s*/\s*\d+|\d+)\s*$")
SLIDE_NUM_RE = re.compile(r"^\s*Slide\s+\d+\s*$", re.IGNORECASE)
EMPTY_BULLET_RE = re.compile(r"^\s*[-•·*o]\s*$")


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def slugify(s: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", s)
    return s.strip("_").upper()


MODULE_SLUGS = {
    "01 - Supplier Management (v176-v178)": "SUPPLIER_MANAGEMENT_V176_V178",
    "02 - Sourcing": "SOURCING",
    "03 - Contract (v182)": "CONTRACT_V182",
    "04 - E-Procurement  (v182)": "E_PROCUREMENT_V182",
    "05 - Invoicing  (v180-182)": "INVOICING_V180_V182",
    "06 - Spend and BI  (v176-v182)": "SPEND_AND_BI_V176_V182",
    "3 - Meet the Ivalua Company (v178)": "MEET_THE_IVALUA_COMPANY_V178",
    "Extranet and PM  (v182)": "EXTRANET_AND_PM_V182",
    "Platform - Admin and Config  (v178-182)": "PLATFORM_ADMIN_AND_CONFIG_V178_V182",
    "Platform - Architecture  (v178-182)": "PLATFORM_ARCHITECTURE_V178_V182",
    "Platform - Integrations (v182)": "PLATFORM_INTEGRATIONS_V182",
    "Soft Skills": "SOFT_SKILLS",
    "_Orientation in the Platform": "ORIENTATION_IN_THE_PLATFORM",
    "_Refresher Badge (v178-v182)": "REFRESHER_BADGE_V178_V182",
}


def module_slug_from_dir(name: str) -> str:
    return MODULE_SLUGS.get(name, slugify(name))


def clean_text(text: str) -> str:
    lines = text.splitlines()
    cleaned = []
    for line in lines:
        if COPYRIGHT_RE.search(line) or PAGE_NUM_RE.match(line) or SLIDE_NUM_RE.match(line) or EMPTY_BULLET_RE.match(line):
            continue
        stripped = line.strip()
        if stripped:
            cleaned.append(stripped)
    return "\n".join(cleaned)


def noise_ratio(original: str, cleaned: str) -> float:
    if not original:
        return 0.0
    removed = max(0, len(original) - len(cleaned))
    return round((removed / len(original)) * 100, 2)


def chunk_text(text: str) -> list[str]:
    if not text:
        return []
    if len(text) <= CHUNK_SIZE:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        if end < len(text):
            break_idx = text.rfind("\n", end - CHUNK_SIZE // 4, end)
            if break_idx == -1:
                break_idx = text.rfind(" ", end - CHUNK_SIZE // 4, end)
            if break_idx != -1 and break_idx > start:
                end = break_idx
        chunks.append(text[start:end].strip())
        next_start = end - CHUNK_OVERLAP
        if next_start <= start:
            next_start = end
        start = next_start
    return chunks


def extract_pdf(fp: Path) -> str:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(str(fp))
        parts = []
        for page in reader.pages:
            try:
                txt = page.extract_text()
                if txt:
                    parts.append(txt)
            except Exception:
                pass
        return "\n".join(parts)
    except Exception as e:
        return f"[PDF extraction error: {e}]"


def extract_docx(fp: Path) -> str:
    try:
        import docx
        return "\n".join(p.text for p in docx.Document(str(fp)).paragraphs if p.text)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_xlsx(fp: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(fp), data_only=True)
        parts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                vals = [str(c.value) for c in row if c.value is not None]
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"


def extract_pptx(fp: Path) -> str:
    try:
        import pptx
        parts = []
        for slide in pptx.Presentation(str(fp)).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract(fp: Path) -> str | None:
    ext = fp.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(fp)
    elif ext == ".txt":
        return fp.read_text(encoding="utf-8", errors="replace")
    elif ext == ".docx":
        return extract_docx(fp)
    elif ext == ".xlsx":
        return extract_xlsx(fp)
    elif ext == ".pptx":
        return extract_pptx(fp)
    return None


def build_conport_items(fp: Path, module_slug: str) -> list[dict]:
    ext = fp.suffix.lower()
    raw = extract(fp)
    if raw is None or not raw.strip():
        return []

    raw = re.sub(r"[\ud800-\udfff]", "", raw)
    original_len = len(raw)
    cleaned = clean_text(raw)
    nr = noise_ratio(raw, cleaned)
    if not cleaned.strip():
        return []

    topic_slug = slugify(fp.stem)[:180]
    chunks = chunk_text(cleaned)
    total = len(chunks)

    items = []
    for idx, chunk in enumerate(chunks):
        key = f"DOC_{module_slug}_{topic_slug}_C{idx}" if total > 1 else f"DOC_{module_slug}_{topic_slug}"
        value = {
            "module": module_slug,
            "topic": fp.stem,
            "file_path": str(fp.relative_to(WORKSPACE)).replace("\\", "/"),
            "content_type": ext.lstrip(".").upper(),
            "chunk_index": idx if total > 1 else None,
            "total_chunks": total if total > 1 else None,
            "noise_ratio": nr,
            "original_char_count": original_len,
            "cleaned_char_count": len(cleaned),
            "char_count": len(chunk),
            "size_bytes": len(json.dumps({"content": chunk}, ensure_ascii=False).encode("utf-8")),
            "content": chunk,
            "is_chunked": total > 1,
            "manual_entry": False,
        }
        items.append({"category": CATEGORY, "key": key, "value": value})
    return items


def scan_docs() -> dict[str, tuple[float, int]]:
    """Return {relative_path: (mtime, size)} for all processable files."""
    files = {}
    for subdir in sorted(DOCS_ROOT.iterdir()):
        if not subdir.is_dir():
            continue
        for fp in sorted(subdir.rglob("*")):
            if not fp.is_file():
                continue
            ext = fp.suffix.lower()
            if ext in SKIP_EXT or ext not in TEXT_EXT:
                continue
            rel = str(fp.relative_to(WORKSPACE)).replace("\\", "/")
            stat = fp.stat()
            files[rel] = (stat.st_mtime, stat.st_size)
    return files


def load_manifest() -> dict:
    if MANIFEST_PATH.exists():
        try:
            with open(MANIFEST_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {"version": 1, "files": {}}


def save_manifest(manifest: dict) -> None:
    MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    tmp = MANIFEST_PATH.with_suffix(".tmp")
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    tmp.replace(MANIFEST_PATH)


def delete_doc_keys(conn: sqlite3.Connection, keys: list[str]) -> int:
    if not keys:
        return 0
    cur = conn.cursor()
    cur.executemany("DELETE FROM custom_data WHERE category = ? AND key = ?", [(CATEGORY, k) for k in keys])
    conn.commit()
    return cur.rowcount


def import_items(conn: sqlite3.Connection, items: list[dict]) -> int:
    if not items:
        return 0
    cur = conn.cursor()
    ts = now_iso()
    rows = [(ts, i["category"], i["key"], json.dumps(i["value"], ensure_ascii=False)) for i in items]
    batch = 200
    total = 0
    for i in range(0, len(rows), batch):
        cur.executemany(
            "INSERT OR REPLACE INTO custom_data (timestamp, category, key, value) VALUES (?, ?, ?, ?)",
            rows[i:i + batch],
        )
        conn.commit()
        total += len(rows[i:i + batch])
    return total


def rebuild_indexes(conn: sqlite3.Connection, all_items: list[dict]) -> int:
    """Rebuild all DOC_INDEX_* entries. Small and fast."""
    by_module: dict[str, list[dict]] = {}
    for item in all_items:
        mod = item["value"]["module"]
        by_module.setdefault(mod, []).append({
            "key": item["key"],
            "topic": item["value"]["topic"],
            "file_path": item["value"]["file_path"],
            "content_type": item["value"]["content_type"],
            "chunk_index": item["value"].get("chunk_index"),
            "total_chunks": item["value"].get("total_chunks"),
        })

    indexes = []
    for mod, topics in by_module.items():
        indexes.append({
            "category": CATEGORY,
            "key": f"DOC_INDEX_{mod}",
            "value": {
                "module_slug": mod,
                "topic_count": len({t["topic"] for t in topics}),
                "chunk_count": len(topics),
                "topics": sorted({t["topic"] for t in topics}),
            },
        })

    indexes.append({
        "category": CATEGORY,
        "key": "DOC_INDEX_ALL",
        "value": {
            "modules": sorted(by_module.keys()),
            "total_modules": len(by_module),
            "total_documents": sum(len({t["topic"] for t in topics}) for topics in by_module.values()),
            "total_chunks": sum(len(topics) for topics in by_module.values()),
        },
    })

    # Delete old indexes
    cur = conn.cursor()
    cur.execute("DELETE FROM custom_data WHERE category = ? AND key LIKE 'DOC_INDEX_%'", (CATEGORY,))
    conn.commit()
    deleted = cur.rowcount

    imported = import_items(conn, indexes)
    return deleted, imported


def main() -> int:
    print("Incremental documentation importer")
    print(f"DB: {DB}")
    print(f"Manifest: {MANIFEST_PATH}")

    conn = sqlite3.connect(DB)
    conn.execute("PRAGMA journal_mode=WAL")

    manifest = load_manifest()
    current_files = scan_docs()
    old_files = manifest.get("files", {})

    # Determine delta
    to_delete = []  # list of keys to remove from DB
    to_process = []  # list of (rel_path, Path) tuples
    unchanged = 0

    for rel, (mtime, size) in current_files.items():
        if rel in old_files:
            old = old_files[rel]
            if old.get("mtime") == mtime and old.get("size") == size:
                unchanged += 1
                continue
        to_process.append(rel)

    removed = [rel for rel in old_files if rel not in current_files]

    print(f"\nScan complete: {len(current_files)} files on disk")
    print(f"  Unchanged: {unchanged}")
    print(f"  Changed/new: {len(to_process)}")
    print(f"  Removed: {len(removed)}")

    # Handle removed files: collect their keys and delete
    if removed:
        removed_keys = []
        for rel in removed:
            removed_keys.extend(old_files[rel].get("keys", []))
        if removed_keys:
            n = delete_doc_keys(conn, removed_keys)
            print(f"  Deleted {n} rows for {len(removed)} removed files")
        for rel in removed:
            del old_files[rel]

    # Handle changed/new files
    processed_items = []
    processed_count = 0
    for rel in to_process:
        fp = WORKSPACE / rel
        module_slug = module_slug_from_dir(fp.parts[len(WORKSPACE.parts)])
        items = build_conport_items(fp, module_slug)
        if items:
            # Delete old keys for this file if it existed
            if rel in old_files:
                old_keys = old_files[rel].get("keys", [])
                if old_keys:
                    delete_doc_keys(conn, old_keys)
            # Insert new items
            import_items(conn, items)
            keys = [i["key"] for i in items]
            old_files[rel] = {
                "mtime": current_files[rel][0],
                "size": current_files[rel][1],
                "module": module_slug,
                "keys": keys,
            }
            processed_items.extend(items)
            processed_count += 1
        else:
            # No extractable content: remove from manifest if present
            if rel in old_files:
                old_keys = old_files[rel].get("keys", [])
                if old_keys:
                    delete_doc_keys(conn, old_keys)
                del old_files[rel]

    print(f"\nProcessed {processed_count} files, {len(processed_items)} chunks")

    # Rebuild indexes (always, fast)
    # Need all items in DB to build correct indexes
    # Re-fetch all doc items from DB
    cur = conn.cursor()
    cur.execute("SELECT key, value FROM custom_data WHERE category = ? AND key NOT LIKE 'DOC_INDEX_%'", (CATEGORY,))
    all_db_items = []
    for key, value in cur.fetchall():
        all_db_items.append({"category": CATEGORY, "key": key, "value": json.loads(value)})

    idx_deleted, idx_imported = rebuild_indexes(conn, all_db_items)
    print(f"Indexes rebuilt: deleted {idx_deleted}, imported {idx_imported}")

    # Save manifest
    manifest["files"] = old_files
    manifest["last_run"] = now_iso()
    save_manifest(manifest)

    # Verify
    cur.execute("SELECT category, COUNT(*) FROM custom_data GROUP BY category ORDER BY category")
    print("\nDB state:")
    for cat, n in cur.fetchall():
        print(f"  {cat}: {n}")

    cur.execute("SELECT COUNT(*), AVG(LENGTH(value)), MAX(LENGTH(value)) FROM custom_data WHERE category = ? AND key NOT LIKE 'DOC_INDEX_%'", (CATEGORY,))
    row = cur.fetchone()
    print(f"\nIVALUA_Documentation stats (excl. indexes):")
    print(f"  Entries: {row[0]}, Avg: {row[1]:.0f} chars, Max: {row[2]} chars")

    conn.close()
    print("\nDone.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
